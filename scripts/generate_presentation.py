from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentations" / "aa-regional-hris-workday-team-agent-overview.pptx"

NAVY = RGBColor(11, 31, 58)
BLUE = RGBColor(30, 78, 121)
TEAL = RGBColor(18, 117, 125)
RED = RGBColor(180, 42, 42)
GOLD = RGBColor(194, 155, 72)
LIGHT = RGBColor(245, 247, 250)
DARK = RGBColor(34, 40, 49)
GRAY = RGBColor(102, 112, 122)
WHITE = RGBColor(255, 255, 255)


def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.7), Inches(0.8))
    title_frame = title_box.text_frame
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE if dark else NAVY

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(11.2), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        paragraph = subtitle_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = subtitle
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT if dark else GRAY


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], color: RGBColor = DARK, font_size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)


def add_kpi_card(slide, left: float, top: float, width: float, height: float, title: str, value: str, subtitle: str, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    frame = shape.text_frame
    frame.clear()

    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE

    paragraph = frame.add_paragraph()
    run = paragraph.add_run()
    run.text = value
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE

    paragraph = frame.add_paragraph()
    run = paragraph.add_run()
    run.text = subtitle
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE


def add_two_column_headers(slide, left_title: str, right_title: str) -> None:
    for x, title in [(0.8, left_title), (6.6, right_title)]:
        box = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(4.8), Inches(0.4))
        frame = box.text_frame
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = title
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = NAVY


def add_banner(slide, text: str, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE


def create_deck() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    blank = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank)
    add_background(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.55), Inches(0.18), Inches(5.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD
    add_title(slide, "AA Regional HRIS WorkDay Team Agent", "Microsoft Foundry + Foundry Agent Service + Microsoft Agent Framework", dark=True)
    add_bullets(
        slide,
        0.95,
        1.9,
        7.2,
        2.6,
        [
            "Solution name: AA Regional HRIS WorkDay Team Agent",
            "One-line impact: Reduce time spent hunting for Workday documentation and integration facts by grounding answers in SharePoint content and a normalized integration catalog.",
            "Team name: AA Regional HRIS WorkDay Team",
            "Team members: To be confirmed",
        ],
        color=WHITE,
        font_size=20,
    )
    add_kpi_card(slide, 8.6, 1.8, 1.8, 1.6, "Lookup time", "5x faster", "Illustrative target for support triage", TEAL)
    add_kpi_card(slide, 10.5, 1.8, 1.8, 1.6, "Coverage", "2 sources", "SharePoint + integration inventory", BLUE)
    add_kpi_card(slide, 8.6, 3.8, 3.7, 1.6, "Trust model", "Cited answers", "Every factual answer points back to a document or spreadsheet row", RED)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, LIGHT)
    add_title(slide, "Business / IT Problem and Value", "Why the AA Workday support team needs this assistant")
    add_bullets(
        slide,
        0.8,
        1.7,
        5.9,
        3.8,
        [
            "Who is impacted: HRIS analysts, Workday integration support, payroll and recruiting operations, and downstream business teams waiting on answers.",
            "Current pain points: Legacy SharePoint content is hard to search, integration facts live in spreadsheets, ownership is inconsistent, and analysts depend on tribal knowledge.",
            "Business / IT impact: Slower triage, higher escalation volume, duplicated investigation work, and longer time to resolve operational questions.",
        ],
        font_size=19,
    )
    add_kpi_card(slide, 7.2, 1.9, 1.7, 1.5, "Time", "30-60 min", "Typical manual lookup across docs and spreadsheets", NAVY)
    add_kpi_card(slide, 9.1, 1.9, 1.7, 1.5, "Target", "< 5 min", "AI-assisted answer with citations", BLUE)
    add_kpi_card(slide, 11.0, 1.9, 1.6, 1.5, "Value", "Lower cost", "Less analyst effort per case", TEAL)
    add_bullets(
        slide,
        7.2,
        3.8,
        5.2,
        2.0,
        [
            "Time savings hypothesis: reclaim analyst time by answering recurring questions immediately.",
            "Cost reduction hypothesis: reduce unnecessary escalations and rework during support investigations.",
            "Customer experience impact: faster answers for internal teams relying on Workday operations.",
        ],
        font_size=18,
    )
    add_banner(slide, "Value comes from faster retrieval, less tribal knowledge dependence, and better answer consistency.", NAVY)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, WHITE)
    add_title(slide, "Vision: What Good Looks Like", "Before vs After the AI assistant")
    add_two_column_headers(slide, "Before", "After")
    left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.2), Inches(3.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(236, 240, 244)
    left_box.line.color.rgb = RGBColor(236, 240, 244)
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(2.0), Inches(5.9), Inches(3.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(228, 243, 242)
    right_box.line.color.rgb = RGBColor(228, 243, 242)
    add_bullets(
        slide,
        1.05,
        2.3,
        4.7,
        3.2,
        [
            "Analyst manually searches SharePoint and spreadsheets.",
            "Support quality depends on who remembers the system history.",
            "Conflicting answers take time to reconcile.",
            "Low-confidence questions become escalations.",
        ],
        font_size=18,
    )
    add_bullets(
        slide,
        6.9,
        2.3,
        5.3,
        3.2,
        [
            "User asks a plain-language question in Teams or a web chat.",
            "AI chooses the right tool path: documents, integration catalog, or both.",
            "Answer includes citations, ownership, and supporting runbooks.",
            "Missing knowledge is logged for backlog improvement instead of being lost.",
        ],
        font_size=18,
    )
    add_banner(slide, "AI changes the experience from manual hunting to guided, grounded operational support.", TEAL)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, LIGHT)
    add_title(slide, "AI Solution Overview", "Inputs → Intelligence → Outputs")
    stages = [
        (0.8, "Inputs", ["SharePoint documentation exports", "Integration spreadsheet / normalized records", "User question from chat or web"]),
        (4.45, "Intelligence", ["Foundry hosted agent orchestration", "Agent Framework tools for retrieval and catalog lookup", "Answer composition with citations and uncertainty handling"]),
        (8.25, "Outputs", ["Grounded natural language answer", "Structured integration facts", "Citations to docs and spreadsheet rows", "Knowledge-gap log for follow-up"]),
    ]
    colors = [NAVY, BLUE, TEAL]
    for index, (left, title, bullets) in enumerate(stages):
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.4), Inches(4.5))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = colors[index]
        label = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.25), Inches(1.45), Inches(1.4), Inches(0.45))
        label.fill.solid()
        label.fill.fore_color.rgb = colors[index]
        label.line.color.rgb = colors[index]
        frame = label.text_frame
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE
        add_bullets(slide, left + 0.22, 2.15, 2.95, 3.7, bullets, font_size=17)
    add_banner(slide, "Design principle: treat SharePoint as unstructured knowledge and the spreadsheet as structured operational data.", NAVY)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, WHITE)
    add_title(slide, "AI Security and Governance", "Security posture for HR-adjacent operational content")
    add_bullets(
        slide,
        0.9,
        1.8,
        5.8,
        4.7,
        [
            "Enterprise authentication only for internal access.",
            "Role-based access aligned to the AA Workday support team and approved stakeholders.",
            "Read-only advisory design for phase one with no Workday write actions.",
            "No training on user conversations by default.",
        ],
        font_size=19,
    )
    add_bullets(
        slide,
        6.8,
        1.8,
        5.6,
        4.7,
        [
            "Every factual answer is expected to carry citations back to source documents or spreadsheet rows.",
            "Redaction and content controls can be added if sensitive fields appear in source data.",
            "Knowledge gaps and low-confidence answers become auditable follow-up items.",
            "Foundry hosting keeps the agent service inside an enterprise-grade operational model.",
        ],
        font_size=19,
    )
    add_banner(slide, "Trust is created by least privilege, grounded answers, and a clear boundary between advice and action.", RED)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, LIGHT)
    add_title(slide, "Demo Walkthrough", "How the user, AI reasoning, and business value connect")
    steps = [
        ("1. User action", "Analyst asks: 'What integrations feed payroll and where is the support documentation?'", NAVY),
        ("2. AI reasoning", "The agent queries the integration catalog for payroll-related records, then retrieves matching SharePoint design and runbook content.", BLUE),
        ("3. Output", "The agent returns the payroll vendor export details, owner team, support contact, cadence, and links the relevant design or runbook.", TEAL),
        ("4. Value", "The analyst gets a grounded answer in one interaction, avoids manual hunting, and can escalate with context if needed.", RED),
    ]
    for index, (title, text, color) in enumerate(steps):
        top = 1.7 + (index * 1.15)
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(top), Inches(2.0), Inches(0.7))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.color.rgb = color
        frame = tag.text_frame
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE

        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(top), Inches(8.9), Inches(0.7))
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        body.line.color.rgb = color
        body_frame = body.text_frame
        paragraph = body_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(17)
        run.font.color.rgb = DARK
    add_banner(slide, "Demo story: one question, two knowledge sources, one cited answer.", TEAL)

    slide = presentation.slides.add_slide(blank)
    add_background(slide, NAVY)
    add_title(slide, "Next Steps", "From MVP to production adoption", dark=True)
    add_bullets(
        slide,
        0.95,
        1.8,
        11.0,
        3.8,
        [
            "Replace sample JSON with real SharePoint export and spreadsheet normalization pipelines.",
            "Configure Foundry project endpoint, model deployment, and environment variables.",
            "Pilot with a small AA Workday support group and capture unanswered questions.",
            "Add stronger retrieval ranking, telemetry, and operational dashboards.",
        ],
        color=WHITE,
        font_size=21,
    )
    add_kpi_card(slide, 0.95, 5.5, 3.4, 1.2, "MVP goal", "Operational assistant", "Grounded answers with citations", BLUE)
    add_kpi_card(slide, 4.55, 5.5, 3.4, 1.2, "Pilot focus", "Analyst productivity", "Faster triage and better consistency", TEAL)
    add_kpi_card(slide, 8.15, 5.5, 3.4, 1.2, "Scale path", "Teams + enterprise data", "Production retrieval and governance", RED)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)


if __name__ == "__main__":
    create_deck()